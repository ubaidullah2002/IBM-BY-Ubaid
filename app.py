# app.py
import streamlit as st
import os
from dotenv import load_dotenv
import logging
import time
from groq import Groq
import docx
import PyPDF2
import pptx
from pathlib import Path
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
import json
import yaml
import io
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import markdown
from PIL import Image
import pdfkit
import imgkit

# Configure page settings for better UI
st.set_page_config(
    page_title="Business Intelligence Suite",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
    .main {
        padding: 0rem 1rem;
    }
    .stApp {
        background-color: #f8f9fa;
    }
    .css-1d391kg {
        padding-top: 1rem;
    }
    .stButton>button {
        background-color: #0066cc;
        color: white;
        border-radius: 5px;
        padding: 0.5rem 1rem;
        border: none;
    }
    .stButton>button:hover {
        background-color: #0052a3;
    }
    .card {
        background-color: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        margin-bottom: 1rem;
    }
    .metric-card {
        text-align: center;
        padding: 1rem;
        background: linear-gradient(135deg, #6699cc 0%, #336699 100%);
        color: white;
        border-radius: 8px;
    }
    /* Story Generator Styles */
    .story-container {
        background-color: white;
        padding: 2rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        margin: 1rem 0;
    }
    
    .resource-section {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 8px;
        margin-top: 2rem;
    }
</style>
""", unsafe_allow_html=True)

# Initialize Groq client
load_dotenv()
client = Groq(api_key=os.getenv('GROQ_API_KEY'))

class BusinessAnalyzer:
    @staticmethod
    def analyze_business_metrics(data):
        try:
            # Analyze key business metrics
            metrics = {
                "revenue_growth": ((data["current_revenue"] - data["previous_revenue"]) / data["previous_revenue"]) * 100,
                "profit_margin": (data["net_profit"] / data["current_revenue"]) * 100,
                "customer_acquisition_cost": data["marketing_spend"] / data["new_customers"],
                "customer_lifetime_value": data["average_order_value"] * data["purchase_frequency"] * data["customer_lifespan"],
                "roi": ((data["net_profit"] - data["total_investment"]) / data["total_investment"]) * 100
            }
            return metrics
        except Exception as e:
            st.error(f"Error analyzing metrics: {str(e)}")
            return None

    @staticmethod
    def generate_business_insights(metrics, business_type):
        prompt = f"""
        Analyze the following business metrics for a {business_type} business and provide strategic insights:
        
        Metrics:
        - Revenue Growth: {metrics['revenue_growth']:.2f}%
        - Profit Margin: {metrics['profit_margin']:.2f}%
        - Customer Acquisition Cost: ${metrics['customer_acquisition_cost']:.2f}
        - Customer Lifetime Value: ${metrics['customer_lifetime_value']:.2f}
        - ROI: {metrics['roi']:.2f}%
        
        Please provide:
        1. Key insights and trends
        2. Specific recommendations for improvement
        3. Potential risks and opportunities
        4. Strategic action items
        """
        
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="mixtral-8x7b-32768",
            temperature=0.3,
            max_tokens=1500
        )
        return response.choices[0].message.content

def process_uploaded_file(uploaded_file):
    file_extension = Path(uploaded_file.name).suffix.lower()
    
    try:
        if file_extension == '.txt':
            return uploaded_file.getvalue().decode()
        elif file_extension == '.docx':
            doc = docx.Document(uploaded_file)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            return '\n'.join([page.extract_text() for page in pdf_reader.pages])
        elif file_extension == '.pptx':
            prs = pptx.Presentation(uploaded_file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    except Exception as e:
        raise Exception(f"Error processing file: {str(e)}")

def process_document(text_content, process_type):
    try:
        prompt = f"""Process the following document content based on {process_type}:
        
        {text_content}
        
        Provide a detailed {process_type} focusing on the main points and key takeaways."""
        
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="mixtral-8x7b-32768",
            temperature=0.3,
            max_tokens=1500
        )
        
        return response.choices[0].message.content
    except Exception as e:
        raise Exception(f"Error processing document: {str(e)}")

def main():
    # Sidebar navigation
    with st.sidebar:
        st.image("https://via.placeholder.com/150x150.png?text=BI+Suite", width=150)
        st.title("Business Suite")
        selected_page = st.radio(
            "Navigate to:",
            ["Business Analytics", "Document Processing", "Contract Generator", 
             "Market Analysis"]
        )
        
        # User profile section in sidebar
        st.sidebar.markdown("---")
        st.sidebar.title("Profile")
        if st.sidebar.button("📋 Dashboard"):
            selected_page = "Business Analytics"
        if st.sidebar.button("📄 Documents"):
            selected_page = "Document Processing"
        st.sidebar.markdown("---")

    if selected_page == "Business Analytics":
        display_business_analytics()
    elif selected_page == "Document Processing":
        display_document_processing()
    elif selected_page == "Contract Generator":
        display_contract_generation()
    elif selected_page == "Market Analysis":
        display_market_analysis()

def display_business_analytics():
    st.title("Business Analytics Dashboard")
    
    # Business Information Input
    with st.expander("Enter Business Metrics", expanded=True):
        col1, col2 = st.columns(2)
        with col1:
            business_type = st.selectbox(
                "Business Type",
                ["E-commerce", "SaaS", "Retail", "Manufacturing", "Services", "Other"]
            )
            current_revenue = st.number_input("Current Revenue ($)", min_value=0.0)
            previous_revenue = st.number_input("Previous Revenue ($)", min_value=0.0)
            net_profit = st.number_input("Net Profit ($)", min_value=0.0)
        
        with col2:
            marketing_spend = st.number_input("Marketing Spend ($)", min_value=0.0)
            new_customers = st.number_input("New Customers", min_value=0)
            average_order_value = st.number_input("Average Order Value ($)", min_value=0.0)
            
        col3, col4 = st.columns(2)
        with col3:
            purchase_frequency = st.number_input("Purchase Frequency (per year)", min_value=0.0)
            customer_lifespan = st.number_input("Customer Lifespan (years)", min_value=0.0)
        with col4:
            total_investment = st.number_input("Total Investment ($)", min_value=0.0)

    if st.button("Analyze Business Metrics"):
        data = {
            "business_type": business_type,
            "current_revenue": current_revenue,
            "previous_revenue": previous_revenue,
            "net_profit": net_profit,
            "marketing_spend": marketing_spend,
            "new_customers": new_customers,
            "average_order_value": average_order_value,
            "purchase_frequency": purchase_frequency,
            "customer_lifespan": customer_lifespan,
            "total_investment": total_investment
        }
        
        metrics = BusinessAnalyzer.analyze_business_metrics(data)
        if metrics:
            # Display metrics in an attractive dashboard layout
            st.subheader("Key Performance Metrics")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("""
                <div class="metric-card">
                    <h3>Revenue Growth</h3>
                    <h2>{:.2f}%</h2>
                </div>
                """.format(metrics["revenue_growth"]), unsafe_allow_html=True)
                
            with col2:
                st.markdown("""
                <div class="metric-card">
                    <h3>Profit Margin</h3>
                    <h2>{:.2f}%</h2>
                </div>
                """.format(metrics["profit_margin"]), unsafe_allow_html=True)
                
            with col3:
                st.markdown("""
                <div class="metric-card">
                    <h3>ROI</h3>
                    <h2>{:.2f}%</h2>
                </div>
                """.format(metrics["roi"]), unsafe_allow_html=True)

            # Generate and display insights
            insights = BusinessAnalyzer.generate_business_insights(metrics, business_type)
            st.subheader("Business Insights")
            st.markdown(f"""
            <div class="card">
                {insights}
            </div>
            """, unsafe_allow_html=True)
            
            # Create visualizations
            fig = go.Figure()
            fig.add_trace(go.Indicator(
                mode="gauge+number",
                value=metrics["profit_margin"],
                title={'text': "Profit Margin"},
                gauge={'axis': {'range': [0, 100]},
                       'bar': {'color': "#336699"}}
            ))
            st.plotly_chart(fig)

def display_document_processing():
    st.title("Smart Document Processing")
    
    # File upload section with enhanced UI
    st.markdown("""
    <div class="card">
        <h3>Upload Documents</h3>
        <p>Support for PDF, Word, PowerPoint, and Text files</p>
    </div>
    """, unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader(
        "",
        type=['txt', 'docx', 'pdf', 'pptx'],
        help="Drag and drop your files here"
    )
    
    if uploaded_file:
        try:
            text_content = process_uploaded_file(uploaded_file)
            st.success(f"Successfully processed {uploaded_file.name}")
            
            # Process options
            process_type = st.selectbox(
                "Select Processing Type",
                ["Smart Summary", "Key Points Extraction", "Action Items", "Full Analysis"]
            )
            
            if st.button("Process Document"):
                with st.spinner("Processing document..."):
                    summary = process_document(text_content, process_type)
                    st.markdown(f"""
                    <div class="card">
                        <h3>Analysis Results</h3>
                        {summary}
                    </div>
                    """, unsafe_allow_html=True)
        except Exception as e:
            st.error(str(e))

def display_contract_generation():
    st.title("AI Contract Generator")
    
    # Initialize contract generator
    contract_gen = ContractGenerator(client)
    
    # Contract configuration
    st.markdown("""
    <div class="card">
        <h3>Contract Setup</h3>
    </div>
    """, unsafe_allow_html=True)
    
    contract_type = st.selectbox(
        "Contract Type",
        ["Service Agreement", "Employment Contract", "NDA", "Custom Contract"]
    )
    
    # Dynamic form based on contract type
    requirements = []
    if contract_type == "Service Agreement":
        with st.expander("Contract Details", expanded=True):
            service_type = st.text_input("Service Type")
            payment_terms = st.text_input("Payment Terms")
            duration = st.text_input("Contract Duration")
            custom_requirements = st.text_area(
                "Additional Requirements (Optional)",
                help="Enter any additional terms or requirements for your contract"
            )
            requirements.extend([
                f"Service Type: {service_type}",
                f"Payment Terms: {payment_terms}",
                f"Duration: {duration}",
                f"Additional Requirements: {custom_requirements}"
            ])
    elif contract_type == "Custom Contract":
        with st.expander("Contract Details", expanded=True):
            custom_requirements = st.text_area(
                "Specify your contract requirements",
                help="Enter specific details about the contract you want to generate. Include any important terms, conditions, or special requirements."
            )
            requirements = [custom_requirements]
    
    if st.button("Generate Contract"):
        with st.spinner("Generating contract..."):
            try:
                # Join requirements into a single string
                requirements_text = "\n".join(filter(None, requirements))
                
                # Generate contract content
                contract_content = contract_gen.generate_contract_template(
                    contract_type, requirements_text
                )
                
                # Display contract preview
                st.markdown("""
                <div class="card">
                    <h3>Contract Preview</h3>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown(contract_content)
                
                # Add download options
                col1, col2 = st.columns(2)
                with col1:
                    # Download as PDF
                    pdf_data = contract_gen.create_pdf(contract_content)
                    st.download_button(
                        "Download as PDF",
                        pdf_data,
                        file_name="contract.pdf",
                        mime="application/pdf"
                    )
                
                with col2:
                    # Download as markdown
                    st.download_button(
                        "Download as Text",
                        contract_content,
                        file_name="contract.md",
                        mime="text/markdown"
                    )
                    
            except Exception as e:
                st.error(f"Error generating contract: {str(e)}")

def display_market_analysis():
    st.title("Market Analysis & Trends")
    
    # Market analysis configuration
    st.markdown("""
    <div class="card">
        <h3>Market Analysis Configuration</h3>
    </div>
    """, unsafe_allow_html=True)
    
    analysis_type = st.selectbox(
        "Select Analysis Type",
        ["Competitor Analysis", "Market Trends", "SWOT Analysis", "Risk Assessment"]
    )
    
    industry = st.selectbox(
        "Industry",
        ["Technology", "Healthcare", "Finance", "Retail", "Manufacturing", "Other"]
    )
    
    # Dynamic form based on analysis type
    if analysis_type == "Competitor Analysis":
        competitors = st.text_area("List main competitors (one per line)")
        market_position = st.selectbox(
            "Your Market Position",
            ["Market Leader", "Strong Competitor", "Growing Player", "New Entrant"]
        )
        
        if st.button("Generate Competitor Analysis"):
            with st.spinner("Analyzing competitors..."):
                analysis = generate_competitor_analysis(competitors, market_position, industry)
                display_analysis_results(analysis)
                
    elif analysis_type == "Market Trends":
        timeframe = st.selectbox("Timeframe", ["Short-term", "Medium-term", "Long-term"])
        focus_areas = st.multiselect(
            "Focus Areas",
            ["Consumer Behavior", "Technology Trends", "Economic Factors", 
             "Regulatory Changes", "Market Size", "Growth Potential"]
        )
        
        if st.button("Analyze Market Trends"):
            with st.spinner("Analyzing market trends..."):
                trends = generate_market_trends(industry, timeframe, focus_areas)
                display_analysis_results(trends)
                
    elif analysis_type == "SWOT Analysis":
        strengths = st.text_area("List your strengths")
        weaknesses = st.text_area("List your weaknesses")
        opportunities = st.text_area("List market opportunities")
        threats = st.text_area("List potential threats")
        
        if st.button("Generate SWOT Analysis"):
            with st.spinner("Generating SWOT analysis..."):
                swot = generate_swot_analysis(
                    strengths, weaknesses, opportunities, threats, industry
                )
                display_analysis_results(swot)
                
    elif analysis_type == "Risk Assessment":
        risk_factors = st.multiselect(
            "Risk Factors to Analyze",
            ["Market Risk", "Financial Risk", "Operational Risk", 
             "Strategic Risk", "Compliance Risk"]
        )
        
        if st.button("Generate Risk Assessment"):
            with st.spinner("Assessing risks..."):
                risks = generate_risk_assessment(risk_factors, industry)
                display_analysis_results(risks)

def generate_competitor_analysis(competitors, market_position, industry):
    prompt = f"""Analyze the competitive landscape for a {market_position} in the {industry} industry.
    
    Competitors:
    {competitors}
    
    Please provide:
    1. Detailed competitor analysis
    2. Market positioning strategy
    3. Competitive advantages and disadvantages
    4. Recommendations for market positioning
    """
    
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="mixtral-8x7b-32768",
        temperature=0.3,
        max_tokens=1500
    )
    return response.choices[0].message.content

def generate_market_trends(industry, timeframe, focus_areas):
    prompt = f"""Analyze market trends for the {industry} industry over a {timeframe} period.
    
    Focus Areas:
    {', '.join(focus_areas)}
    
    Please provide:
    1. Current market trends
    2. Future projections
    3. Impact analysis
    4. Strategic recommendations
    """
    
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="mixtral-8x7b-32768",
        temperature=0.3,
        max_tokens=1500
    )
    return response.choices[0].message.content

def generate_swot_analysis(strengths, weaknesses, opportunities, threats, industry):
    prompt = f"""Perform a SWOT analysis for a company in the {industry} industry.
    
    Strengths:
    {strengths}
    
    Weaknesses:
    {weaknesses}
    
    Opportunities:
    {opportunities}
    
    Threats:
    {threats}
    
    Please provide:
    1. Detailed SWOT analysis
    2. Strategic implications
    3. Recommended actions
    """
    
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="mixtral-8x7b-32768",
        temperature=0.3,
        max_tokens=1500
    )
    return response.choices[0].message.content

def generate_risk_assessment(risk_factors, industry):
    prompt = f"""Assess risks for a company in the {industry} industry.
    
    Risk Factors:
    {', '.join(risk_factors)}
    
    Please provide:
    1. Risk analysis for each factor
    2. Risk mitigation strategies
    3. Monitoring recommendations
    4. Contingency planning
    """
    
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="mixtral-8x7b-32768",
        temperature=0.3,
        max_tokens=1500
    )
    return response.choices[0].message.content

def display_analysis_results(analysis):
    st.markdown("""
    <div class="card">
        <h3>Analysis Results</h3>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown(analysis)
    
    # Download options
    st.download_button(
        "Download Analysis",
        analysis,
        file_name="market_analysis.md",
        mime="text/markdown"
    )

class ContractGenerator:
    def __init__(self, client):
        self.client = client

    def generate_contract_template(self, contract_type, requirements):
        prompt = f"""Generate a professional {contract_type} contract with the following requirements:
        {requirements}
        
        Please include all standard legal sections including:
        1. Parties involved
        2. Terms and conditions
        3. Payment terms (if applicable)
        4. Duration
        5. Termination clauses
        6. Governing law
        7. Signature blocks
        
        Format in proper legal contract style with clear sections and numbering."""

        response = self.client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="mixtral-8x7b-32768",
            temperature=0.3,
            max_tokens=2000
        )
        return response.choices[0].message.content

    def create_pdf(self, content):
        try:
            # Create a buffer for the PDF
            buffer = io.BytesIO()
            
            # Create the PDF object
            pdf = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter
            
            # Set font and size
            pdf.setFont("Helvetica", 12)
            
            # Split content into lines
            y = height - 50  # Start from top with margin
            lines = content.split('\n')
            
            for line in lines:
                # Wrap text to fit page width
                words = line.split()
                current_line = []
                
                for word in words:
                    current_line.append(word)
                    # Check if current line width exceeds page width
                    if pdf.stringWidth(' '.join(current_line)) > width - 100:
                        current_line.pop()
                        # Draw the line
                        pdf.drawString(50, y, ' '.join(current_line))
                        y -= 20  # Move down for next line
                        current_line = [word]
                        
                        # Check if we need a new page
                        if y < 50:
                            pdf.showPage()
                            y = height - 50
                            pdf.setFont("Helvetica", 12)
                
                # Draw remaining words in the line
                if current_line:
                    pdf.drawString(50, y, ' '.join(current_line))
                    y -= 20
            
            pdf.save()
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            st.error(f"Error creating PDF: {str(e)}")
            return None

    def create_image(self, content):
        try:
            # Convert markdown to HTML first
            html_content = markdown.markdown(content)
            
            # Create image using imgkit with error handling
            options = {
                'quiet': '',
                'encoding': 'UTF-8'
            }
            img = imgkit.from_string(html_content, False, options=options)
            return img
        except Exception as e:
            st.error(f"Error creating image: {str(e)}")
            return None

if __name__ == "__main__":
    main()